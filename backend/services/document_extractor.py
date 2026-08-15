import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import io
import logging
import os
import base64
import requests

logger = logging.getLogger(__name__)

def ocr_pdf_pages(doc) -> str:
    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        return ""
    parts = []
    for i in range(len(doc)):
        pix = doc[i].get_pixmap(dpi=72)
        img_bytes = pix.tobytes("jpeg", 70) 
        encoded = base64.b64encode(img_bytes).decode('utf-8')
        parts.append({"inlineData": {"mimeType": "image/jpeg", "data": encoded}})
        
    parts.append({"text": "Extract all text from these document pages sequentially. Maintain the original formatting, headings, and lists. Do not summarize, just perform pure OCR."})
    
    payload = {
        "contents": [{"role": "user", "parts": parts}],
        "generationConfig": {"temperature": 0, "maxOutputTokens": 8192}
    }
    
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-3.5-flash:generateContent?key={api_key}"
    try:
        res = requests.post(url, headers={"Content-Type": "application/json"}, json=payload, timeout=90)
        if res.status_code == 200:
            return res.json()["candidates"][0]["content"]["parts"][0]["text"].strip()
    except Exception as e:
        logger.error(f"OCR failed: {e}")
    return ""

def extract_text_from_document(file_bytes: bytes, filename: str) -> str:
    ext = filename.lower().split('.')[-1]
    if ext == 'pdf':
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        text = ""
        for page in doc:
            text += page.get_text()
        text = text.strip()
        
        # If average chars per page is very low, it's a scanned PDF!
        if len(text) < (len(doc) * 50):
            logger.info("PDF appears to be scanned. Using Gemini multi-page OCR fallback...")
            ocr_text = ocr_pdf_pages(doc)
            if ocr_text:
                text = (text + "\n" + ocr_text).strip()
        return text
        
    elif ext == 'docx':
        doc = Document(io.BytesIO(file_bytes))
        return "\n".join([para.text for para in doc.paragraphs]).strip()
        
    elif ext == 'pptx':
        prs = Presentation(io.BytesIO(file_bytes))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text).strip()
        
    elif ext == 'txt':
        return file_bytes.decode('utf-8', errors='ignore').strip()
        
    else:
        raise ValueError(f"Unsupported file format: {ext}")
